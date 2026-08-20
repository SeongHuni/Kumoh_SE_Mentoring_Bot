# 2차 발표 PPTX 생성기
#   python scripts/build-ppt.py
#
# 서사 원칙: "성과"가 아니라 "시행착오"를 중심에 둔다.
# 각 슬라이드는 [무엇을 시도했나 → 어떻게 됐나 → 왜 그랬나]의 흐름을 따른다.
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

OUT = os.path.join("docs", "학과챗봇_RAG_2차발표.pptx")

# ---- 팔레트 ----
ACCENT = RGBColor(0x0E, 0x7C, 0x86)   # 딥 틸 — 강조
ACCENT_L = RGBColor(0xDC, 0xED, 0xEE)  # 연한 틸 — 배경 강조
INK = RGBColor(0x16, 0x20, 0x1F)       # 본문
MUTED = RGBColor(0x66, 0x74, 0x6F)     # 보조 텍스트
LINE = RGBColor(0xD6, 0xDE, 0xDE)      # 구분선
BG = RGBColor(0xFF, 0xFF, 0xFF)
PANEL = RGBColor(0xF2, 0xF5, 0xF5)
POS = RGBColor(0x2E, 0x7D, 0x5B)       # 개선
NEG = RGBColor(0xAA, 0x40, 0x38)       # 하락 / 문제
WARN_BG = RGBColor(0xF6, 0xEB, 0xD2)
WARN_FG = RGBColor(0x8A, 0x62, 0x12)

FONT = "맑은 고딕"
MONO = "Consolas"

W, H = Inches(13.333), Inches(7.5)
M = Inches(0.75)          # 좌우 여백
CONTENT_W = W - 2 * M


def set_font(run, size, bold=False, color=INK, name=FONT):
    """python-pptx는 라틴 폰트만 설정하므로 동아시아 폰트도 함께 지정한다."""
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = name
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set("typeface", name)


def textbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    return tf


def para(tf, text, size, bold=False, color=INK, space_after=6, first=False,
         name=FONT, align=PP_ALIGN.LEFT, line=None):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    if line:
        p.line_spacing = line
    run = p.add_run()
    run.text = text
    set_font(run, size, bold, color, name)
    return p


def rect(slide, x, y, w, h, fill=PANEL, line_color=None):
    from pptx.enum.shapes import MSO_SHAPE
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line_color:
        sh.line.color.rgb = line_color
        sh.line.width = Pt(1)
    else:
        sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def new_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])  # 빈 레이아웃


def header(slide, eyebrow, title, n=None):
    """모든 내용 슬라이드 상단 공통 — 라벨 + 제목 + 밑줄."""
    tf = textbox(slide, M, Inches(0.45), CONTENT_W, Inches(0.3))
    para(tf, eyebrow, 13, True, ACCENT, 0, first=True)

    tf2 = textbox(slide, M, Inches(0.82), CONTENT_W, Inches(0.9))
    para(tf2, title, 30, True, INK, 0, first=True, line=1.15)

    ln = rect(slide, M, Inches(1.72), Inches(1.1), Pt(3), ACCENT)
    if n is not None:
        tf3 = textbox(slide, W - M - Inches(0.6), H - Inches(0.55), Inches(0.6), Inches(0.3))
        para(tf3, f"{n:02d}", 11, False, LINE, 0, first=True, name=MONO, align=PP_ALIGN.RIGHT)


def bullets(slide, x, y, w, items, size=17, gap=11):
    """items: (텍스트, 굵게) 튜플 목록. '•' 마커를 직접 그린다."""
    tf = textbox(slide, x, y, w, Inches(4))
    for i, (text, bold) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        p.line_spacing = 1.3
        r1 = p.add_run(); r1.text = "•  "
        set_font(r1, size, False, ACCENT)
        r2 = p.add_run(); r2.text = text
        set_font(r2, size, bold, INK if bold else MUTED)
    return tf


def table(slide, x, y, w, rows, col_w=None, size=14, header_bg=ACCENT,
          highlight=None, align_right_from=1):
    """rows[0]은 헤더. highlight: 강조할 데이터 행 인덱스 목록(0-based, 헤더 제외)."""
    nrow, ncol = len(rows), len(rows[0])
    h = Inches(0.42) * nrow
    shape = slide.shapes.add_table(nrow, ncol, x, y, w, h)
    tbl = shape.table

    if col_w:
        total = sum(col_w)
        for i, cw in enumerate(col_w):
            tbl.columns[i].width = Emu(int(w * cw / total))

    for r in range(nrow):
        tbl.rows[r].height = Inches(0.42)
        for c in range(ncol):
            cell = tbl.cell(r, c)
            cell.margin_left = Inches(0.12)
            cell.margin_right = Inches(0.12)
            cell.margin_top = Inches(0.04)
            cell.margin_bottom = Inches(0.04)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

            txt = str(rows[r][c])
            cell.fill.solid()
            if r == 0:
                cell.fill.fore_color.rgb = header_bg
                color, bold = RGBColor(0xFF, 0xFF, 0xFF), True
            else:
                is_hl = highlight and (r - 1) in highlight
                cell.fill.fore_color.rgb = ACCENT_L if is_hl else BG
                bold = is_hl or c == 0
                color = INK
                if txt.startswith("+"):
                    color = POS
                elif txt.startswith("−") or txt.startswith("-"):
                    color = NEG

            tfc = cell.text_frame
            tfc.word_wrap = True
            p = tfc.paragraphs[0]
            p.alignment = PP_ALIGN.RIGHT if (c >= align_right_from and r > 0) else PP_ALIGN.LEFT
            if r == 0 and c >= align_right_from:
                p.alignment = PP_ALIGN.RIGHT
            run = p.add_run()
            run.text = txt
            # 숫자 칸은 고정폭으로 자릿수를 맞춘다
            use_mono = c >= align_right_from and r > 0 and any(ch.isdigit() for ch in txt)
            set_font(run, size, bold, color, MONO if use_mono else FONT)
    return tbl


def big_stat(slide, x, y, w, value, caption, color=ACCENT, vsize=54):
    """핵심 수치 하나를 크게."""
    h = Inches(1.5)
    rect(slide, x, y, w, h, ACCENT_L)
    tf = textbox(slide, x + Inches(0.3), y + Inches(0.18), w - Inches(0.6), Inches(0.8))
    para(tf, value, vsize, True, color, 2, first=True, name=MONO)
    tf2 = textbox(slide, x + Inches(0.3), y + Inches(1.0), w - Inches(0.6), Inches(0.4))
    para(tf2, caption, 14, False, INK, 0, first=True)


def callout(slide, x, y, w, label, text, bg=WARN_BG, fg=WARN_FG, h=Inches(1.15)):
    rect(slide, x, y, w, h, bg)
    rect(slide, x, y, Pt(4), h, fg)
    tf = textbox(slide, x + Inches(0.25), y + Inches(0.14), w - Inches(0.5), Inches(0.3))
    para(tf, label, 12, True, fg, 4, first=True)
    tf2 = textbox(slide, x + Inches(0.25), y + Inches(0.48), w - Inches(0.5), h - Inches(0.6))
    para(tf2, text, 15, False, INK, 0, first=True, line=1.25)


def build():
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H

    # ================= 1. 표지 =================
    s = new_slide(prs)
    rect(s, Inches(0), Inches(0), W, H, BG)
    rect(s, Inches(0), Inches(0), Inches(0.22), H, ACCENT)

    tf = textbox(s, Inches(1.2), Inches(2.0), Inches(10.5), Inches(0.4))
    para(tf, "소프트웨어전공 SIG · 2차 발표", 16, True, ACCENT, 0, first=True)

    tf = textbox(s, Inches(1.2), Inches(2.6), Inches(11), Inches(1.8))
    para(tf, "잘 된 것보다\n안 된 것을 공유합니다", 44, True, INK, 0, first=True, line=1.2)

    tf = textbox(s, Inches(1.2), Inches(4.6), Inches(10.5), Inches(1.0))
    para(tf, "RAG 성능을 올리려고 당연하게 적용했던 두 가지 기법이", 18, False, MUTED, 4, first=True)
    para(tf, "왜 우리 데이터에서는 통하지 않았는가", 18, False, MUTED, 0)

    rect(s, Inches(1.2), Inches(5.9), Inches(3.6), Pt(2), LINE)
    tf = textbox(s, Inches(1.2), Inches(6.1), Inches(10), Inches(0.4))
    para(tf, "Overlap · 리랭킹(Reranking)", 15, True, INK, 0, first=True)

    # ================= 2. 오늘 할 이야기 =================
    s = new_slide(prs)
    header(s, "OVERVIEW", "오늘 할 이야기", 2)

    items = [
        ("지난 발표에서 청킹 방식을 정했습니다 — 문서를 500 토큰씩 자르기", False),
        ("이번엔 성능을 더 올리려고 두 가지 기법을 추가로 적용했습니다", False),
        ("결과는 둘 다 실패였습니다. 오늘은 그 과정을 공유합니다", True),
    ]
    bullets(s, M, Inches(2.1), CONTENT_W, items, size=18)

    y = Inches(3.5)
    cards = [
        ("시도 ①", "Overlap", "청크를 겹쳐서 자르기", "효과 없음"),
        ("시도 ②", "리랭킹", "검색 결과 순위 다시 매기기", "효과 없음"),
        ("다음", "프롬프트 전략", "대화 맥락 처리하기", "설계 완료"),
    ]
    cw = (CONTENT_W - Inches(0.5)) / 3
    for i, (tag, name, desc, res) in enumerate(cards):
        x = M + i * (cw + Inches(0.25))
        rect(s, x, y, cw, Inches(2.2), PANEL)
        rect(s, x, y, cw, Pt(4), ACCENT if i < 2 else MUTED)
        tf = textbox(s, x + Inches(0.28), y + Inches(0.3), cw - Inches(0.56), Inches(0.3))
        para(tf, tag, 12, True, ACCENT if i < 2 else MUTED, 6, first=True)
        para(tf, name, 24, True, INK, 8)
        para(tf, desc, 14, False, MUTED, 12)
        para(tf, res, 15, True, NEG if i < 2 else POS, 0)

    # ================= 3. 평가 방법 =================
    s = new_slide(prs)
    header(s, "배경", "먼저, 어떻게 평가했는지", 3)

    tf = textbox(s, M, Inches(2.05), Inches(6.4), Inches(0.4))
    para(tf, "질문 80개와 정답을 미리 만들어두고, 검색이 정답을 찾아내는지 셌습니다.",
         17, False, INK, 0, first=True, line=1.3)

    rect(s, M, Inches(2.7), Inches(6.4), Inches(2.5), PANEL)
    tf = textbox(s, M + Inches(0.3), Inches(2.9), Inches(5.8), Inches(2.1))
    for i, (k, v) in enumerate([
        ('"question"', '"2025-2학기 수강지도 상담 기간은?"'),
        ('"gold_doc_ids"', '["kumoh-notice-44578"]'),
        ('"expected_evidence"', '["7. 28.(월) 09:00 ~ ..."]'),
        ('"answer_checkpoints"', '["시작일과 종료일", "시간"]'),
    ]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(9)
        r1 = p.add_run(); r1.text = k + ": "
        set_font(r1, 13, True, ACCENT, MONO)
        r2 = p.add_run(); r2.text = v
        set_font(r2, 13, False, INK, MONO)

    tf = textbox(s, M + Inches(0.3), Inches(4.75), Inches(5.8), Inches(0.4))
    para(tf, "질문 · 정답 문서 · 근거 · 채점 기준", 12, False, MUTED, 0, first=True)

    x2 = M + Inches(6.9)
    w2 = CONTENT_W - Inches(6.9)
    tf = textbox(s, x2, Inches(2.05), w2, Inches(0.4))
    para(tf, "핵심 지표", 13, True, ACCENT, 8, first=True)

    rect(s, x2, Inches(2.5), w2, Inches(1.3), ACCENT_L)
    tf = textbox(s, x2 + Inches(0.3), Inches(2.7), w2 - Inches(0.6), Inches(1.0))
    para(tf, "Recall@3", 26, True, ACCENT, 6, first=True, name=MONO)
    para(tf, "정답 문서가 검색 상위 3개 안에\n들어온 질문의 비율", 14, False, INK, 0, line=1.25)

    tf = textbox(s, x2, Inches(4.1), w2, Inches(1.6))
    para(tf, "예) Recall@3 = 0.60 이면", 14, True, INK, 6, first=True)
    para(tf, "질문 100개 중 60개는 상위 3개 안에\n정답이 있었다는 뜻", 14, False, MUTED, 0, line=1.3)

    # ================= 4. 시행착오 ① =================
    s = new_slide(prs)
    header(s, "시행착오 ①", "Overlap을 넣었는데, 아무 일도 일어나지 않았다", 4)

    tf = textbox(s, M, Inches(2.05), CONTENT_W, Inches(0.5))
    para(tf, "Overlap = 청크를 자를 때 앞부분을 조금 겹쳐서 자르는 방법. "
             "문장이 중간에 끊기는 걸 막으려고 대부분의 RAG가 기본으로 씁니다.",
         16, False, MUTED, 0, first=True, line=1.3)

    rows = [
        ["겹치는 양", "0 (안 겹침)", "50", "100", "150"],
        ["Recall@3", "0.875", "0.850", "0.838", "0.863"],
        ["Recall@5", "0.912", "0.950", "0.900", "0.938"],
        ["청크 개수", "1,474", "1,503", "1,531", "1,581"],
    ]
    table(s, M, Inches(2.85), Inches(7.4), rows, col_w=[2.0, 1.5, 1.2, 1.2, 1.2], size=15,
          highlight=[1])

    x2 = M + Inches(7.8)
    w2 = CONTENT_W - Inches(7.8)
    big_stat(s, x2, Inches(2.85), w2, "68 / 80", "질문 80개 중 결과가\n전혀 바뀌지 않은 개수")

    callout(s, M, Inches(4.9), CONTENT_W,
            "이상한 점",
            "겹치는 양을 50 → 100 → 150으로 늘렸는데 점수가 오르락내리락합니다. "
            "진짜 효과가 있다면 이렇게 들쭉날쭉하지 않습니다.")

    # ================= 5. 원인 ① =================
    s = new_slide(prs)
    header(s, "원인을 찾아보니 ①", "문서 절반 이상은 애초에 잘리지도 않았다", 5)

    rows = [
        ["문서가 나뉜 개수", "문서 수", "비율"],
        ["1개 (안 나뉨)", "496", "58.4%"],
        ["2개", "208", "24.5%"],
        ["3개", "79", "9.3%"],
        ["4개 이상", "67", "7.9%"],
    ]
    table(s, M, Inches(2.2), Inches(6.2), rows, col_w=[2.4, 1.2, 1.2], size=15, highlight=[0])

    x2 = M + Inches(6.7)
    w2 = CONTENT_W - Inches(6.7)
    big_stat(s, x2, Inches(2.2), w2, "58.4%", "문서 하나가 청크 하나 — 자를 일이 없었음", vsize=50)

    tf = textbox(s, x2, Inches(3.95), w2, Inches(1.2))
    para(tf, "우리 문서 길이 중간값", 14, True, INK, 4, first=True)
    para(tf, "445 토큰", 30, True, ACCENT, 4, name=MONO)
    para(tf, "자르는 기준(500 토큰)보다 짧다", 14, False, MUTED, 0)

    callout(s, M, Inches(5.4), CONTENT_W,
            "그래서",
            "문서가 안 나뉘면 겹칠 앞부분 자체가 없습니다. "
            "전체의 58%에서는 Overlap이 물리적으로 적용될 수가 없었습니다.")

    # ================= 6. 원인 ② =================
    s = new_slide(prs)
    header(s, "원인을 찾아보니 ②", "공지사항은 소설이 아니다", 6)

    left_w = Inches(6.0)
    tf = textbox(s, M, Inches(2.1), left_w, Inches(2.6))
    para(tf, "Overlap이 필요한 문서", 15, True, NEG, 8, first=True)
    para(tf, "소설 · 논문처럼 앞 내용을 알아야\n뒤 내용이 이해되는 글", 16, False, MUTED, 20, line=1.3)
    para(tf, "우리 문서", 15, True, ACCENT, 8)
    para(tf, "\"신청 기간\", \"제출 서류\", \"유의사항\"처럼\n문단마다 내용이 따로 노는 공지사항",
         16, False, MUTED, 0, line=1.3)

    x2 = M + Inches(6.4)
    w2 = CONTENT_W - Inches(6.4)
    big_stat(s, x2, Inches(2.1), w2, "89%", "청크가 나뉜 지점이\n문장 단위로 깔끔하게 끊김", vsize=52)

    tf = textbox(s, x2, Inches(3.85), w2, Inches(1.4))
    para(tf, "우리는 자를 때 문단 → 목록 → 문장 순서로", 14, False, MUTED, 4, first=True, line=1.3)
    para(tf, "자연스러운 경계에서만 자릅니다.", 14, False, MUTED, 0, line=1.3)

    callout(s, M, Inches(5.35), CONTENT_W,
            "결론",
            "문장이 중간에 끊기는 경우가 10%뿐이었습니다. Overlap이 고쳐줄 문제가 거의 없었던 겁니다.")

    # ================= 7. 검증 =================
    s = new_slide(prs)
    header(s, "그럼 확인해보자", "문서가 길어지면 Overlap이 살아날까?", 7)

    tf = textbox(s, M, Inches(2.05), CONTENT_W, Inches(0.4))
    para(tf, "같은 데이터를 문서 길이별로 나눠서 다시 계산해봤습니다.",
         17, False, MUTED, 0, first=True)

    rows = [
        ["정답 문서 길이", "질문 수", "Overlap 없음", "Overlap 적용", "차이"],
        ["500 토큰 이하", "142", "0.634", "0.613", "−0.021"],
        ["500 ~ 1,000", "90", "0.589", "0.578", "−0.011"],
        ["1,000 ~ 1,500", "40", "0.525", "0.625", "+0.100"],
        ["1,500 이상", "28", "0.500", "0.571", "+0.071"],
    ]
    table(s, M, Inches(2.65), Inches(8.6), rows,
          col_w=[2.2, 1.1, 1.5, 1.5, 1.1], size=15, highlight=[2, 3])

    x2 = M + Inches(9.0)
    w2 = CONTENT_W - Inches(9.0)
    rect(s, x2, Inches(2.65), w2, Inches(1.05), RGBColor(0xDF, 0xEE, 0xE7))
    tf = textbox(s, x2 + Inches(0.25), Inches(2.8), w2 - Inches(0.5), Inches(0.8))
    para(tf, "긴 문서에서는", 13, True, POS, 3, first=True)
    para(tf, "효과가 나타났다", 20, True, POS, 0)

    rect(s, x2, Inches(3.85), w2, Inches(1.05), PANEL)
    tf = textbox(s, x2 + Inches(0.25), Inches(3.98), w2 - Inches(0.5), Inches(0.85))
    para(tf, "긴 문서를 묻는 질문은", 13, False, MUTED, 2, first=True)
    para(tf, "23%뿐", 24, True, ACCENT, 0, name=MONO)

    callout(s, M, Inches(5.35), CONTENT_W,
            "그래서 이렇게 정리했습니다",
            "Overlap은 쓸모없는 기법이 아닙니다. 효과가 나는 조건이 있고, "
            "우리 데이터의 23%만 거기에 해당해서 전체적으로는 상쇄됐습니다.")

    # ================= 8. 시행착오 ② =================
    s = new_slide(prs)
    header(s, "시행착오 ②", '"500 토큰에서만 본 거 아니야?" — 다시 해봤다', 8)

    tf = textbox(s, M, Inches(2.05), CONTENT_W, Inches(0.5))
    para(tf, "처음엔 500 토큰 하나에서만 실험했습니다. 지적을 받고 "
             "자르는 크기와 자르는 방식을 바꿔가며 다시 측정했습니다.",
         16, False, MUTED, 0, first=True, line=1.3)

    rows = [
        ["실험 설정", "Overlap 없음", "Overlap 적용", "차이"],
        ["300 토큰씩 자르기", "0.617", "0.610", "−0.007"],
        ["500 토큰씩 자르기", "0.593", "0.600", "+0.007"],
        ["1000 토큰씩 자르기", "0.623", "0.620", "−0.003"],
        ["다른 방식(recursive) 500", "0.597", "0.587", "−0.010"],
    ]
    table(s, M, Inches(2.9), Inches(8.2), rows, col_w=[2.6, 1.4, 1.4, 1.1], size=15)

    x2 = M + Inches(8.6)
    w2 = CONTENT_W - Inches(8.6)
    big_stat(s, x2, Inches(2.9), w2, "0.01", "어떤 설정에서도\n차이는 이 정도 미만", vsize=46)

    callout(s, M, Inches(5.2), CONTENT_W,
            "추가로 확인한 것",
            "혹시 Overlap이 제대로 안 걸린 건 아닌지도 봤습니다. "
            "recursive 방식에서는 모든 경계에 100% 적용됐는데도 결과는 같았습니다.")

    # ================= 9. Overlap 결론 =================
    s = new_slide(prs)
    header(s, "정리", "Overlap — 쓰지 않기로 했습니다", 9)

    half = (CONTENT_W - Inches(0.4)) / 2
    rect(s, M, Inches(2.2), half, Inches(2.4), PANEL)
    tf = textbox(s, M + Inches(0.3), Inches(2.45), half - Inches(0.6), Inches(2.0))
    para(tf, "좋아진 것", 15, True, MUTED, 12, first=True)
    para(tf, "확인되지 않음", 24, True, MUTED, 12)
    para(tf, "· 어떤 설정에서도 차이 0.01 미만\n· 통계적으로도 차이 없음\n· 80문항 중 68~72개는 변화 없음",
         14, False, MUTED, 0, line=1.5)

    x2 = M + half + Inches(0.4)
    rect(s, x2, Inches(2.2), half, Inches(2.4), RGBColor(0xF5, 0xE2, 0xE0))
    tf = textbox(s, x2 + Inches(0.3), Inches(2.45), half - Inches(0.6), Inches(2.0))
    para(tf, "나빠진 것", 15, True, NEG, 12, first=True)
    para(tf, "확인됨", 24, True, NEG, 12)
    para(tf, "· 청크 개수 2~7% 증가\n· 저장·임베딩 비용 증가\n· 비슷한 청크가 검색 상위 자리를 차지",
         14, False, INK, 0, line=1.5)

    rect(s, M, Inches(4.9), CONTENT_W, Inches(1.15), ACCENT_L)
    tf = textbox(s, M + Inches(0.4), Inches(5.15), CONTENT_W - Inches(0.8), Inches(0.7))
    para(tf, "좋아진 건 확인 못 했고 나빠진 건 확인했으니, 굳이 복잡하게 만들 이유가 없다고 판단했습니다.",
         18, True, INK, 0, first=True, line=1.25)

    tf = textbox(s, M, Inches(6.3), CONTENT_W, Inches(0.4))
    para(tf, "※ 우리 데이터 기준입니다. 소설이나 스캔 문서처럼 문맥이 이어지는 자료라면 결과가 다를 수 있습니다.",
         13, False, MUTED, 0, first=True)

    # ================= 10. 시행착오 ③ =================
    s = new_slide(prs)
    header(s, "시행착오 ③", "리랭킹도 붙여봤지만 마찬가지였다", 10)

    tf = textbox(s, M, Inches(2.05), CONTENT_W, Inches(0.5))
    para(tf, "리랭킹 = 검색으로 후보 50개를 넉넉히 가져온 뒤, 더 똑똑한 모델로 순위를 다시 매기는 방법.",
         16, False, MUTED, 0, first=True, line=1.3)

    rows = [
        ["방식", "Recall@3", "Recall@5", "걸린 시간"],
        ["리랭킹 안 함 (기준)", "0.593", "0.740", "0 ms"],
        ["BGE 모델", "0.520", "0.690", "1,736 ms"],
        ["Qwen3 모델", "0.630", "0.727", "2,222 ms"],
    ]
    table(s, M, Inches(2.8), Inches(7.6), rows, col_w=[2.4, 1.2, 1.2, 1.4], size=16)

    x2 = M + Inches(8.0)
    w2 = CONTENT_W - Inches(8.0)
    big_stat(s, x2, Inches(2.8), w2, "500배", "검색은 3ms인데\n리랭킹은 2.2초가 걸립니다", vsize=46, color=NEG)

    callout(s, M, Inches(4.85), CONTENT_W,
            "판단",
            "점수는 거의 그대로인데 응답이 2초 넘게 느려집니다. "
            "지금 단계에서는 그만한 값어치가 없다고 봤습니다.",
            bg=RGBColor(0xF5, 0xE2, 0xE0), fg=NEG)

    # ================= 11. 리랭킹 원인 =================
    s = new_slide(prs)
    header(s, "왜 그럴까 (추정)", "우리 문서들이 서로 너무 비슷하다", 11)

    tf = textbox(s, M, Inches(2.05), CONTENT_W, Inches(0.5))
    para(tf, "리랭커는 '관련 있는 문서'와 '관련 없는 문서'를 가려내는 데 강합니다. "
             "그런데 우리 후보는 전부 어느 정도 관련 있는 비슷한 문서였습니다.",
         16, False, MUTED, 0, first=True, line=1.3)

    y = Inches(2.85)
    cw = (CONTENT_W - Inches(0.6)) / 3
    facts = [
        ("76%", "데이터의 대부분이\n학과 게시판 한 곳에서 나옴"),
        ("30%", "연도·학기만 다른\n똑같은 제목의 공지"),
        ("55%", "검색 상위 5개의 점수가\n거의 붙어 있는 질문 비율"),
    ]
    for i, (v, c) in enumerate(facts):
        x = M + i * (cw + Inches(0.3))
        rect(s, x, y, cw, Inches(1.75), PANEL)
        tf = textbox(s, x + Inches(0.3), y + Inches(0.25), cw - Inches(0.6), Inches(1.3))
        para(tf, v, 40, True, ACCENT, 8, first=True, name=MONO)
        para(tf, c, 14, False, MUTED, 0, line=1.3)

    tf = textbox(s, M, Inches(4.85), CONTENT_W, Inches(0.5))
    para(tf, "학과 게시판 3년치를 모은 데이터라, 한 분야 안의 비슷한 문서만 있습니다. "
             "리랭커 입장에서는 우열을 가릴 단서가 부족했던 것으로 보입니다.",
         16, False, INK, 0, first=True, line=1.3)

    callout(s, M, Inches(5.6), CONTENT_W,
            "다만 — 아직 확실하지 않습니다",
            "이건 저희 추정이고, 검증이 끝난 결론은 아닙니다. 다음 장에서 설명하겠습니다.",
            h=Inches(1.0))

    # ================= 12. 아직 못 푼 것 =================
    s = new_slide(prs)
    header(s, "솔직하게", "아직 해결하지 못한 문제들", 12)

    tf = textbox(s, M, Inches(2.05), CONTENT_W, Inches(0.4))
    para(tf, "리랭킹 쪽은 Overlap만큼 검증이 끝나지 않았습니다.", 17, False, MUTED, 0, first=True)

    issues = [
        ("설정 실수 가능성", "리랭커에 입력 길이를 지정하지 않아서, 청크 뒷부분이 잘린 채 "
                            "들어갔을 수 있습니다. BGE 성적이 유독 나쁜 이유가 이것일 수 있습니다."),
        ("모델 사용법 확인 필요", "Qwen3 모델을 원래 쓰는 방식과 다르게 불러왔을 가능성이 있습니다."),
        ("결과 파일 유실", "300문항 실행 결과가 화면 캡처로만 남아 있습니다. 다시 돌려서 파일로 남겨야 합니다."),
    ]
    y = Inches(2.7)
    for i, (t, d) in enumerate(issues):
        rect(s, M, y, CONTENT_W, Inches(1.15), PANEL)
        rect(s, M, y, Pt(4), Inches(1.15), NEG)
        tf = textbox(s, M + Inches(0.3), y + Inches(0.16), CONTENT_W - Inches(0.6), Inches(0.3))
        para(tf, f"{i+1}. {t}", 16, True, NEG, 5, first=True)
        tf2 = textbox(s, M + Inches(0.3), y + Inches(0.53), CONTENT_W - Inches(0.6), Inches(0.55))
        para(tf2, d, 14, False, INK, 0, first=True, line=1.25)
        y += Inches(1.32)

    rect(s, M, Inches(6.75), CONTENT_W, Pt(3), ACCENT)
    tf = textbox(s, M, Inches(6.9), CONTENT_W, Inches(0.4))
    para(tf, "그래서 리랭킹 결과는 \"현재 설정 기준\"으로 봐주시면 됩니다.",
         15, True, INK, 0, first=True)

    # ================= 13. 다음 단계 =================
    s = new_slide(prs)
    header(s, "다음 단계", "이제 프롬프트 전략을 만듭니다", 13)

    tf = textbox(s, M, Inches(2.05), CONTENT_W, Inches(0.5))
    para(tf, "지금 챗봇은 질문 하나만 받습니다. \"그럼 승인 안 되면?\" 같은 이어지는 질문을 못 알아듣습니다.",
         17, False, MUTED, 0, first=True, line=1.3)

    y = Inches(2.8)
    steps = [
        ("1", "이어지는 질문을 완전한 질문으로 바꾼다",
         '"그럼 승인 안 되면?" → "2026-2학기 수강지도 상담 미승인 시 제한"'),
        ("2", "그 질문으로 검색하되, 공지와 후기를 구분한다",
         "일정·신청은 공식 공지만 / 강의 후기는 에브리타임만"),
        ("3", "검색된 자료만으로 답하고 출처를 붙인다",
         "무엇을 묻는지 모르면 추측하지 않고 되묻는다"),
    ]
    for n, t, d in steps:
        rect(s, M, y, Inches(0.55), Inches(1.05), ACCENT)
        tf = textbox(s, M, y + Inches(0.28), Inches(0.55), Inches(0.5))
        para(tf, n, 22, True, RGBColor(0xFF, 0xFF, 0xFF), 0, first=True,
             name=MONO, align=PP_ALIGN.CENTER)
        tf2 = textbox(s, M + Inches(0.85), y + Inches(0.14), CONTENT_W - Inches(1.0), Inches(0.9))
        para(tf2, t, 17, True, INK, 5, first=True)
        para(tf2, d, 14, False, MUTED, 0, line=1.25)
        y += Inches(1.25)

    rect(s, M, Inches(6.6), CONTENT_W, Inches(0.75), ACCENT_L)
    tf = textbox(s, M + Inches(0.35), Inches(6.78), CONTENT_W - Inches(0.7), Inches(0.4))
    para(tf, "핵심 — 이전 대화는 '무엇을 찾을지' 정할 때만 쓰고, 답변의 근거로는 쓰지 않습니다.",
         16, True, INK, 0, first=True)

    # ================= 14. 마무리 =================
    s = new_slide(prs)
    header(s, "마무리", "이번에 배운 것", 14)

    rows = [
        ["", "결정", "이유"],
        ["청킹 500 토큰", "그대로 유지", "우리 문서 길이(중간값 445)와 잘 맞음"],
        ["Overlap", "쓰지 않음", "좋아진 건 확인 못 하고 비용만 늘어남"],
        ["리랭킹", "지금은 보류", "점수는 그대로인데 2.2초 느려짐"],
        ["프롬프트 전략", "다음에 구현", "이어지는 질문 처리가 필요함"],
    ]
    table(s, M, Inches(2.15), CONTENT_W, rows, col_w=[2.2, 1.6, 5.0], size=16, align_right_from=99)

    y = Inches(4.55)
    lessons = [
        "관행이라고 그냥 쓰지 않고, 넣어보고 재보고 근거를 갖고 뺐습니다",
        "안 되는 이유를 찾다 보니 우리 데이터의 성격을 훨씬 잘 알게 됐습니다",
        "아직 검증이 덜 끝난 부분은 덜 끝났다고 말하는 게 맞다고 생각합니다",
    ]
    tf = textbox(s, M, y, CONTENT_W, Inches(2.0))
    for i, t in enumerate(lessons):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(14)
        p.line_spacing = 1.3
        r1 = p.add_run(); r1.text = "✓  "
        set_font(r1, 17, True, ACCENT)
        r2 = p.add_run(); r2.text = t
        set_font(r2, 17, False, INK)

    os.makedirs("docs", exist_ok=True)
    prs.save(OUT)
    print(f"생성 완료: {OUT}")
    print(f"슬라이드 {len(prs.slides.__iter__.__self__._sldIdLst)}장")


if __name__ == "__main__":
    build()
